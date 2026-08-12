#!/usr/bin/env python3
"""
pptx_generator.py — Generador de Presentaciones PowerPoint (.pptx) para el Dimensionador Estratégico.

Genera un deck ejecutivo de 16:9 con diseño corporativo Iris StartUp Lab (Design_2.md)
aplicando paleta morado/dorado, tarjetas de KPI, badges de veredicto a color y 
fichas ejecutorias por idea.

Uso:
    python scripts/pptx_generator.py [--data report_data.json] [--output deck_priorizacion.pptx]
"""

import sys
import json
import argparse
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# --- Constantes y Paleta de Colores (Design_2.md) ---
COLOR_PURPLE_DARK  = RGBColor(0x24, 0x1B, 0x33)  # #241B33 - Headers y Dark Theme
COLOR_PURPLE_MED   = RGBColor(0x3D, 0x27, 0x66)  # #3D2766 - Accent Purple
COLOR_PURPLE_LIGHT = RGBColor(0xF7, 0xF3, 0xFC)  # #F7F3FC - Cards & Zebra
COLOR_PURPLE_CARD  = RGBColor(0xED, 0xE6, 0xF7)  # #EDE6F7 - KPI Cards
COLOR_GOLD         = RGBColor(0xE8, 0xB9, 0x3E)  # #E8B93E - Highlights
COLOR_GOLD_DARK    = RGBColor(0xD4, 0xA7, 0x3E)  # #D4A73E - Gold Text/Border
COLOR_WHITE        = RGBColor(0xFF, 0xFF, 0xFF)  # #FFFFFF - Text on Dark

# Semáforo de Veredictos
COLOR_GREEN_TEXT   = RGBColor(0x15, 0x80, 0x3D)  # PROTOTIPAR
COLOR_GREEN_BG     = RGBColor(0xDC, 0xFC, 0xE7)

COLOR_GOLD_BG      = RGBColor(0xFE, 0xF9, 0xC3)  # VALIDAR MÁS
COLOR_AMBER_TEXT   = RGBColor(0x92, 0x40, 0x0E)

COLOR_RED_TEXT     = RGBColor(0xB8, 0x4A, 0x3D)  # DESCARTAR
COLOR_RED_BG       = RGBColor(0xFE, 0xE2, 0xE2)

FONT_TITLE = "Sora"
FONT_BODY  = "Inter"


def create_solid_shape(slide, shape_type, left, top, width, height, bg_color, border_color=None):
    """Crea una forma geométrica con color sólido de fondo y borde opcional."""
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, margin=0):
    """Crea una caja de texto sin bordes ni rellenos predeterminados."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    return tf


def add_header(slide, title_text, category_text="IRIS STARTUP LAB — DIMENSIONADOR ESTRATÉGICO"):
    """Agrega un encabezado corporativo estándar para diapositivas claras."""
    # Eyebrow
    tf_eyebrow = add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.3))
    p0 = tf_eyebrow.paragraphs[0]
    p0.text = category_text.upper()
    p0.font.name = FONT_TITLE
    p0.font.size = Pt(9)
    p0.font.bold = True
    p0.font.color.rgb = COLOR_GOLD_DARK

    # Title
    tf_title = add_text_box(slide, Inches(0.8), Inches(0.7), Inches(11.7), Inches(0.6))
    p1 = tf_title.paragraphs[0]
    p1.text = title_text
    p1.font.name = FONT_TITLE
    p1.font.size = Pt(22)
    p1.font.bold = True
    p1.font.color.rgb = COLOR_PURPLE_DARK


def get_verdict_colors(verdict):
    """Devuelve (bg_color, text_color) según el veredicto."""
    v = str(verdict).upper()
    if "PROTOTIPAR" in v:
        return COLOR_GREEN_BG, COLOR_GREEN_TEXT
    elif "VALIDAR" in v:
        return COLOR_GOLD_BG, COLOR_AMBER_TEXT
    else:
        return COLOR_RED_BG, COLOR_RED_TEXT


def generate_pptx_deck(data, output_path="deck_priorizacion_dimensionador.pptx"):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]  # Layout en blanco

    meta = data.get("meta", {})
    ideas = data.get("ideas", [])

    # ==========================================
    # SLIDE 1: PORTADA EJECUTIVA (DARK THEME)
    # ==========================================
    slide1 = prs.slides.add_slide(blank_layout)
    create_solid_shape(slide1, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5), COLOR_PURPLE_DARK)

    # Accent Gold Line Top
    create_solid_shape(slide1, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.15), COLOR_GOLD)

    # Eyebrow
    tf = add_text_box(slide1, Inches(1.0), Inches(1.5), Inches(11.333), Inches(0.4))
    p = tf.paragraphs[0]
    p.text = "IRIS STARTUP LAB — DIMENSIONADOR ESTRATÉGICO v3.0"
    p.font.name = FONT_TITLE
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = COLOR_GOLD

    # Title
    tf = add_text_box(slide1, Inches(1.0), Inches(2.0), Inches(11.333), Inches(1.2))
    p = tf.paragraphs[0]
    p.text = "Deck de Priorización de Ideas de Innovación"
    p.font.name = FONT_TITLE
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE

    # Subtitle / Context Box
    create_solid_shape(slide1, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(3.6), Inches(11.333), Inches(2.8), COLOR_PURPLE_MED)
    tf_ctx = add_text_box(slide1, Inches(1.3), Inches(3.8), Inches(10.733), Inches(2.4))

    p = tf_ctx.paragraphs[0]
    p.text = "CONTEXTO Y CALIBRACIÓN DEL PORTAFOLIO"
    p.font.name = FONT_TITLE
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = COLOR_GOLD

    context_items = [
        ("Objetivo Estratégico", meta.get("objetivoEstrategico", "📈 Incrementar mercado")),
        ("Etapa del Negocio", meta.get("etapaNegocio", "Growth / Corporativo")),
        ("Sector / Vertical", meta.get("sector", "Fintech B2B / SaaS")),
        ("Geografía Objetivo", meta.get("geografia", "México / LATAM")),
        ("Recursos de Prototipado", meta.get("recursosPrototipado", "$50,000 USD / 90 días"))
    ]

    for label, val in context_items:
        p = tf_ctx.add_paragraph()
        run1 = p.add_run()
        run1.text = f"• {label}: "
        run1.font.bold = True
        run1.font.size = Pt(13)
        run1.font.color.rgb = COLOR_WHITE
        run1.font.name = FONT_BODY

        run2 = p.add_run()
        run2.text = str(val)
        run2.font.size = Pt(13)
        run2.font.color.rgb = COLOR_PURPLE_LIGHT
        run2.font.name = FONT_BODY

    # Footer Slide 1
    tf_ft = add_text_box(slide1, Inches(1.0), Inches(6.7), Inches(11.333), Inches(0.4))
    p = tf_ft.paragraphs[0]
    p.text = "Documento Confidencial · Evaluación Financiera y Estratégica Tipo McKinsey / VC"
    p.font.name = FONT_BODY
    p.font.size = Pt(10)
    p.font.color.rgb = COLOR_PURPLE_LIGHT

    # ==========================================
    # SLIDE 2: TABLA RESUMEN DE PRIORIZACIÓN
    # ==========================================
    slide2 = prs.slides.add_slide(blank_layout)
    create_solid_shape(slide2, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5), COLOR_PURPLE_LIGHT)
    add_header(slide2, "Matriz Resumen de Priorización de Ideas", "PARTE 1 — MATRIZ DE DECISIÓN")

    # Table creation
    rows = len(ideas) + 1
    cols = 8
    left = Inches(0.8)
    top = Inches(1.5)
    width = Inches(11.733)
    height = Inches(0.5 * rows)

    table_shape = slide2.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table

    # Column Widths
    table.columns[0].width = Inches(0.6)   # #
    table.columns[1].width = Inches(3.2)   # Idea
    table.columns[2].width = Inches(1.3)   # Modelo
    table.columns[3].width = Inches(1.2)   # Score
    table.columns[4].width = Inches(1.4)   # TAM
    table.columns[5].width = Inches(1.4)   # SOM 3y
    table.columns[6].width = Inches(1.1)   # CLV:CAC
    table.columns[7].width = Inches(1.533) # Veredicto

    headers = ["#", "Idea de Negocio", "Modelo", "Score /25", "TAM Global", "SOM 3 Años", "CLV:CAC", "Veredicto"]
    for col_idx, h_text in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR_PURPLE_DARK
        p = cell.text_frame.paragraphs[0]
        p.text = h_text
        p.font.name = FONT_TITLE
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = COLOR_WHITE
        p.alignment = PP_ALIGN.CENTER if col_idx not in (1, 2) else PP_ALIGN.LEFT

    for row_idx, idea in enumerate(ideas, start=1):
        bg_col = COLOR_WHITE if row_idx % 2 == 1 else COLOR_PURPLE_CARD
        row_data = [
            str(idea.get("rank", row_idx)),
            idea.get("name", "Idea"),
            idea.get("model", "N/A"),
            f"{idea.get('score', 0)}/25",
            idea.get("tam", "N/A"),
            idea.get("som3y", "N/A"),
            idea.get("clvCacAdjusted", "N/A"),
            idea.get("verdict", "PENDIENTE")
        ]

        for col_idx, val in enumerate(row_data):
            cell = table.cell(row_idx, col_idx)
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg_col
            p = cell.text_frame.paragraphs[0]
            p.text = val
            p.font.name = FONT_BODY
            p.font.size = Pt(10)
            
            if col_idx == 7: # Verdict column
                v_bg, v_txt = get_verdict_colors(val)
                cell.fill.fore_color.rgb = v_bg
                p.font.bold = True
                p.font.color.rgb = v_txt
                p.alignment = PP_ALIGN.CENTER
            elif col_idx in (0, 3, 4, 5, 6):
                p.alignment = PP_ALIGN.CENTER
                p.font.color.rgb = COLOR_PURPLE_DARK
                if col_idx == 3:
                    p.font.bold = True

    # ==========================================
    # SLIDE 3+: FICHAS EJECUTIVAS POR IDEA
    # ==========================================
    for idea in ideas:
        slide = prs.slides.add_slide(blank_layout)
        create_solid_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5), COLOR_PURPLE_LIGHT)

        # Header Bar
        rank_num = idea.get("rank", 1)
        name_str = idea.get("name", "Idea de Negocio")
        model_str = idea.get("model", "Modelo N/A")
        verdict_str = idea.get("verdict", "PROTOTIPAR")

        add_header(slide, f"Idea #{rank_num}: {name_str}", f"FICHA EJECUTIVA · MODELO: {model_str}")

        # Verdict Badge Right Top
        v_bg, v_txt = get_verdict_colors(verdict_str)
        create_solid_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(10.2), Inches(0.5), Inches(2.3), Inches(0.6), v_bg)
        tf_v = add_text_box(slide, Inches(10.2), Inches(0.55), Inches(2.3), Inches(0.5))
        p = tf_v.paragraphs[0]
        p.text = verdict_str
        p.alignment = PP_ALIGN.CENTER
        p.font.name = FONT_TITLE
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = v_txt

        # 3 Column Layout
        col_w = Inches(3.7)
        gap = Inches(0.3)
        top_pos = Inches(1.5)
        h_pos = Inches(5.5)

        # ----------------------------------------------------
        # COLUMNA 1: MARKET SIZING & FUENTES
        # ----------------------------------------------------
        left1 = Inches(0.8)
        create_solid_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left1, top_pos, col_w, h_pos, COLOR_WHITE, COLOR_PURPLE_CARD)
        
        # Col 1 Header
        create_solid_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left1, top_pos, col_w, Inches(0.5), COLOR_PURPLE_DARK)
        tf_h1 = add_text_box(slide, left1 + Inches(0.1), top_pos + Inches(0.08), col_w - Inches(0.2), Inches(0.4))
        p = tf_h1.paragraphs[0]
        p.text = "1. Dimensionamiento & Fuentes"
        p.font.name = FONT_TITLE
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = COLOR_WHITE

        tf_c1 = add_text_box(slide, left1 + Inches(0.2), top_pos + Inches(0.6), col_w - Inches(0.4), h_pos - Inches(0.7))
        
        # Metrics list
        metrics1 = [
            ("TAM (Global/Reg):", idea.get("tam", "N/A")),
            ("SAM (Accesible):", idea.get("sam", "N/A")),
            ("SOM 1 Año:", idea.get("som1y", "N/A")),
            ("SOM 3 Años:", idea.get("som3y", "N/A")),
            ("SOM 5 Años:", idea.get("som5y", "N/A"))
        ]
        for lbl, val in metrics1:
            p = tf_c1.add_paragraph() if tf_c1.paragraphs[0].text else tf_c1.paragraphs[0]
            r1 = p.add_run(); r1.text = f"{lbl} "; r1.font.bold = True; r1.font.size = Pt(10); r1.font.name = FONT_BODY
            r2 = p.add_run(); r2.text = str(val); r2.font.bold = True; r2.font.color.rgb = COLOR_GOLD_DARK; r2.font.size = Pt(10); r2.font.name = FONT_BODY

        p_src = tf_c1.add_paragraph()
        p_src.text = "\nFuentes Verificadas / Supuestos:"
        p_src.font.bold = True
        p_src.font.size = Pt(10)
        p_src.font.name = FONT_TITLE
        p_src.font.color.rgb = COLOR_PURPLE_DARK

        sources = idea.get("sourcesList", [
            {"name": "Banxico / INEGI", "verified": True},
            {"name": "Estimación interna", "verified": False}
        ])[:3]
        for src in sources:
            p = tf_c1.add_paragraph()
            flag = "[✓ Verificada]" if src.get("verified") else "[⚠️ Estimación]"
            p.text = f"• {flag} {src.get('name')}"
            p.font.size = Pt(9)
            p.font.name = FONT_BODY
            p.font.color.rgb = COLOR_GREEN_TEXT if src.get("verified") else COLOR_AMBER_TEXT

        # ----------------------------------------------------
        # COLUMNA 2: UNIT ECONOMICS & CROSS-SELLING
        # ----------------------------------------------------
        left2 = left1 + col_w + gap
        create_solid_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left2, top_pos, col_w, h_pos, COLOR_WHITE, COLOR_PURPLE_CARD)

        create_solid_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left2, top_pos, col_w, Inches(0.5), COLOR_PURPLE_DARK)
        tf_h2 = add_text_box(slide, left2 + Inches(0.1), top_pos + Inches(0.08), col_w - Inches(0.2), Inches(0.4))
        p = tf_h2.paragraphs[0]
        p.text = "2. Unit Economics & Cross-Sell"
        p.font.name = FONT_TITLE
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = COLOR_WHITE

        tf_c2 = add_text_box(slide, left2 + Inches(0.2), top_pos + Inches(0.6), col_w - Inches(0.4), h_pos - Inches(0.7))
        
        ue_list = idea.get("unitEconomics", [{}])
        ue_item = ue_list[0] if len(ue_list) > 0 else {}

        metrics2 = [
            ("CLV Base:", ue_item.get("clvBase", "N/A")),
            ("Boost Cross-Sell:", ue_item.get("crossSellBoost", "N/A")),
            ("CLV Ajustado:", ue_item.get("clvAdjusted", "N/A")),
            ("CAC Estimado:", ue_item.get("cac", "N/A")),
            ("Ratio CLV:CAC:", idea.get("clvCacAdjusted", "N/A")),
            ("Payback Period:", ue_item.get("payback", "N/A"))
        ]
        for lbl, val in metrics2:
            p = tf_c2.add_paragraph() if tf_c2.paragraphs[0].text else tf_c2.paragraphs[0]
            r1 = p.add_run(); r1.text = f"{lbl} "; r1.font.bold = True; r1.font.size = Pt(10); r1.font.name = FONT_BODY
            r2 = p.add_run(); r2.text = str(val); r2.font.bold = True; r2.font.color.rgb = COLOR_PURPLE_DARK; r2.font.size = Pt(10); r2.font.name = FONT_BODY

        p_bp = tf_c2.add_paragraph()
        p_bp.text = "\nBuyer Personas Aplicables:"
        p_bp.font.bold = True
        p_bp.font.size = Pt(10)
        p_bp.font.name = FONT_TITLE

        bps = idea.get("buyerPersonas", ["BP-1: Segmento Principal"])
        for bp in bps[:2]:
            p = tf_c2.add_paragraph()
            p.text = f"• {bp}"
            p.font.size = Pt(9)
            p.font.name = FONT_BODY

        # ----------------------------------------------------
        # COLUMNA 3: SCORE /25, RIESGOS & VEREDICTO
        # ----------------------------------------------------
        left3 = left2 + col_w + gap
        create_solid_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left3, top_pos, col_w, h_pos, COLOR_WHITE, COLOR_PURPLE_CARD)

        create_solid_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left3, top_pos, col_w, Inches(0.5), COLOR_PURPLE_DARK)
        tf_h3 = add_text_box(slide, left3 + Inches(0.1), top_pos + Inches(0.08), col_w - Inches(0.2), Inches(0.4))
        p = tf_h3.paragraphs[0]
        p.text = "3. Score /25 & Riesgos Críticos"
        p.font.name = FONT_TITLE
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = COLOR_WHITE

        tf_c3 = add_text_box(slide, left3 + Inches(0.2), top_pos + Inches(0.6), col_w - Inches(0.4), h_pos - Inches(0.7))

        # Score total
        p = tf_c3.paragraphs[0]
        p.text = f"SCORE FINAL: {idea.get('score', 0)} / 25 PTS"
        p.font.name = FONT_TITLE
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = COLOR_GOLD_DARK

        # Score breakdown by criterion
        score_breakdown = idea.get("scoreBreakdown", [])
        if score_breakdown:
            p_sb_h = tf_c3.add_paragraph()
            p_sb_h.text = "\nDesglose por Criterio:"
            p_sb_h.font.bold = True
            p_sb_h.font.size = Pt(9)
            p_sb_h.font.name = FONT_TITLE
            p_sb_h.font.color.rgb = COLOR_PURPLE_DARK

            for criterion_item in score_breakdown[:5]:
                c_name = criterion_item.get("criterion", "Criterio")
                c_pts  = criterion_item.get("points", 0)
                # Visual bar using Unicode blocks (max 5 chars = 5 pts)
                filled = "█" * int(c_pts)
                empty  = "░" * (5 - int(c_pts))
                p_c = tf_c3.add_paragraph()
                p_c.text = f"{filled}{empty} {c_pts}/5  {c_name}"
                p_c.font.size = Pt(8)
                p_c.font.name = FONT_BODY
                p_c.font.color.rgb = COLOR_PURPLE_DARK

        p_r = tf_c3.add_paragraph()
        p_r.text = "\nRiesgos Críticos:"
        p_r.font.bold = True
        p_r.font.size = Pt(9)
        p_r.font.name = FONT_TITLE

        risks_raw = idea.get("risks", ["Adopción inicial de mercado", "Regulación sectorial"])
        risks = []
        for r in risks_raw[:2]:
            if isinstance(r, dict):
                risks.append(r.get("risk", str(r)))
            else:
                risks.append(str(r))

        for r in risks:
            p = tf_c3.add_paragraph()
            p.text = f"! {r}"
            p.font.size = Pt(8)
            p.font.name = FONT_BODY
            p.font.color.rgb = COLOR_RED_TEXT

        p_v = tf_c3.add_paragraph()
        p_v.text = f"\nVeredicto: {verdict_str}"
        p_v.font.bold = True
        p_v.font.size = Pt(9)
        p_v.font.name = FONT_TITLE
        p_v.font.color.rgb = v_txt

        desc_v = idea.get("verdictReason", "Aprobado para prototipado rápido en ciclo de 90 días.")
        p_desc = tf_c3.add_paragraph()
        p_desc.text = desc_v[:160] + ("..." if len(desc_v) > 160 else "")
        p_desc.font.size = Pt(8)
        p_desc.font.name = FONT_BODY

    # ==========================================
    # SLIDE FINAL: HOJA DE RUTA (DARK THEME)
    # ==========================================
    slide_final = prs.slides.add_slide(blank_layout)
    create_solid_shape(slide_final, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5), COLOR_PURPLE_DARK)

    # Accent Gold Line Top
    create_solid_shape(slide_final, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.15), COLOR_GOLD)

    tf = add_text_box(slide_final, Inches(1.0), Inches(1.0), Inches(11.333), Inches(0.4))
    p = tf.paragraphs[0]
    p.text = "HOJA DE RUTA Y PRÓXIMOS PASOS"
    p.font.name = FONT_TITLE
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = COLOR_GOLD

    tf = add_text_box(slide_final, Inches(1.0), Inches(1.5), Inches(11.333), Inches(0.8))
    p = tf.paragraphs[0]
    p.text = "Plan de Prototipado y Validación a 90 Días"
    p.font.name = FONT_TITLE
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE

    # 3 Phase Cards
    phases = [
        ("FASE 1 (Días 1–30)", "Validación de Supuestos Críticos", "Entrevistas en profundidad con BP prioritario y landing page de prueba para validar intención de pago (CPA/CPL)."),
        ("FASE 2 (Días 31–60)", "Desarrollo de MVP Funcional", "Construcción del prototipo de baja/media fidelidad enfocado en la propuesta de valor diferenciadora."),
        ("FASE 3 (Días 61–90)", "Pruebas Piloto y Unit Economics", "Lanzamiento controlado con cohortes iniciales para medir tasa de conversión real y retención.")
    ]

    card_w = Inches(3.6)
    card_gap = Inches(0.26)
    card_top = Inches(2.6)
    card_h = Inches(3.8)

    for idx, (f_title, f_sub, f_desc) in enumerate(phases):
        c_left = Inches(1.0) + idx * (card_w + card_gap)
        create_solid_shape(slide_final, MSO_SHAPE.ROUNDED_RECTANGLE, c_left, card_top, card_w, card_h, COLOR_PURPLE_MED)

        tf_ph = add_text_box(slide_final, c_left + Inches(0.2), card_top + Inches(0.3), card_w - Inches(0.4), card_h - Inches(0.6))
        
        p = tf_ph.paragraphs[0]
        p.text = f_title
        p.font.name = FONT_TITLE
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = COLOR_GOLD

        p = tf_ph.add_paragraph()
        p.text = f_sub
        p.font.name = FONT_TITLE
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = COLOR_WHITE

        p = tf_ph.add_paragraph()
        p.text = f"\n{f_desc}"
        p.font.name = FONT_BODY
        p.font.size = Pt(11)
        p.font.color.rgb = COLOR_PURPLE_LIGHT

    prs.save(output_path)
    print(f"[OK] Presentacion PowerPoint generada exitosamente en: {output_path}")


def get_default_data():
    return {
        "meta": {
            "objetivoEstrategico": "📈 Incrementar mercado y capturar nuevos segmentos",
            "etapaNegocio": "Growth / Corporativo",
            "sector": "Fintech B2B / SaaS",
            "geografia": "México / LATAM",
            "recursosPrototipado": "$50,000 USD / 90 días"
        },
        "ideas": [
            {
                "id": "idea-1",
                "rank": 1,
                "name": "Plataforma de Factoring Digital B2B",
                "model": "B2B SaaS + Take Rate",
                "score": 23,
                "verdict": "PROTOTIPAR",
                "tam": "$1.2B USD",
                "sam": "$180M USD",
                "som1y": "$1.5M USD",
                "som3y": "$8.2M USD",
                "som5y": "$22.0M USD",
                "clvCacAdjusted": "4.8:1",
                "buyerPersonas": ["BP-1: PyMEs Proveedoras", "BP-2: Grandes Compradores"],
                "sourcesList": [
                    {"name": "Banxico Informes Sectoriales", "verified": True},
                    {"name": "INEGI Estadística Industrial", "verified": True}
                ],
                "unitEconomics": [{
                    "clvBase": "$3,200 USD",
                    "crossSellBoost": "+$850 USD",
                    "clvAdjusted": "$4,050 USD",
                    "cac": "$840 USD",
                    "payback": "4.5 meses"
                }],
                "risks": ["Riesgo de crédito de pagadores", "Integración con ERPs"],
                "verdictReason": "Excelente ratio CLV:CAC de 4.8:1 con alto potencial de penetración en PyMEs de México."
            },
            {
                "id": "idea-2",
                "rank": 2,
                "name": "Score de Crédito Alternativo con IA",
                "model": "API Transaccional",
                "score": 18,
                "verdict": "VALIDAR MÁS",
                "tam": "$850M USD",
                "sam": "$95M USD",
                "som1y": "$450K USD",
                "som3y": "$3.1M USD",
                "som5y": "$9.5M USD",
                "clvCacAdjusted": "3.2:1",
                "buyerPersonas": ["BP-1: Sofomes / Neobancos"],
                "sourcesList": [
                    {"name": "CNBV Reportes Financieros", "verified": True},
                    {"name": "Estimación de volumen transaccional", "verified": False}
                ],
                "unitEconomics": [{
                    "clvBase": "$1,800 USD",
                    "crossSellBoost": "+$300 USD",
                    "clvAdjusted": "$2,100 USD",
                    "cac": "$650 USD",
                    "payback": "6.0 meses"
                }],
                "risks": ["Regulación Ley Fintech", "Precisión del algoritmo de scoring"],
                "verdictReason": "Requiere pruebas de precisión del algoritmo antes de iniciar prototipado comercial."
            }
        ]
    }


def main():
    parser = argparse.ArgumentParser(description="Generador de Presentaciones PowerPoint (.pptx) para el Dimensionador Estratégico.")
    parser.add_argument("--data", help="Ruta al archivo JSON con los datos del reporte")
    parser.add_argument("--output", default="deck_priorizacion_dimensionador.pptx", help="Ruta de salida del archivo .pptx")
    args = parser.parse_args()

    if args.data:
        try:
            with open(args.data, "r", encoding="utf-8") as f:
                data = json.load(f)
        except Exception as e:
            print(f"[⚠️] Error al leer {args.data}: {e}. Usando datos de prueba por defecto.")
            data = get_default_data()
    else:
        data = get_default_data()

    generate_pptx_deck(data, args.output)


if __name__ == "__main__":
    main()
